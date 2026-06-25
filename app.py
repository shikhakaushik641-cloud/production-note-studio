import streamlit as st
import pypdf
from pptx import Presentation
import zipfile
import io
import re
import os

# ── API Keys — from env vars; sidebar input overrides at runtime ──────────────
GEMINI_KEYS = [k.strip() for k in os.environ.get("GEMINI_API_KEYS", "").split(",") if k.strip()]
CLAUDE_KEYS = [k.strip() for k in os.environ.get("CLAUDE_API_KEYS", "").split(",") if k.strip()]
OPENAI_KEYS = [k.strip() for k in os.environ.get("OPENAI_API_KEYS", "").split(",") if k.strip()]

# ═══════════════════════════════════════════════════════════════════════════════
#  CSS  —  Modern Scientific Design System
# ═══════════════════════════════════════════════════════════════════════════════
CSS = """
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800&family=Lora:ital,wght@0,400;0,600;1,400&family=JetBrains+Mono:wght@400;500&display=swap');

/* ── Design Tokens ── */
:root {
    /* Scientific palette */
    --indigo:       #1a237e;
    --indigo-mid:   #283593;
    --indigo-lite:  #e8eaf6;
    --blue:         #0d47a1;
    --blue-mid:     #1565c0;
    --blue-lite:    #e3f2fd;
    --emerald:      #2e7d32;
    --emerald-lite: #e8f5e9;
    --amber:        #ff8f00;
    --amber-lite:   #fff8e1;
    --amethyst:     #6a1b9a;
    --amethyst-lite:#f3e5f5;
    --teal:         #00695c;
    --teal-lite:    #e0f2f1;
    --rose:         #c62828;
    --rose-lite:    #ffebee;
    /* Neutrals */
    --ink:          #1a1a2e;
    --ink-muted:    #455a64;
    --surface:      #ffffff;
    --canvas:       #f8f9fc;
    --canvas2:      #f0f2f8;
    --border:       #dde3ef;
    --rule:         #e8eaf0;
    /* Shadows */
    --shadow-sm:    0 1px 3px rgba(26,35,126,0.08), 0 1px 2px rgba(0,0,0,0.06);
    --shadow-md:    0 4px 6px rgba(26,35,126,0.07), 0 2px 4px rgba(0,0,0,0.06);
    --shadow-lg:    0 10px 25px rgba(26,35,126,0.10), 0 4px 10px rgba(0,0,0,0.06);
    /* Typography */
    --fn:           'Inter', system-ui, -apple-system, sans-serif;
    --fn-serif:     'Lora', Georgia, serif;
    --fn-mono:      'JetBrains Mono', 'Fira Code', Consolas, monospace;
    --radius:       12px;
    --radius-sm:    8px;
}

*, *::before, *::after { box-sizing: border-box; margin: 0; padding: 0; }
html { scroll-behavior: smooth; }

body, .stMarkdown, .stTabs [data-baseweb="tab-panel"] {
    background: var(--canvas);
    font-family: var(--fn);
    font-size: 16px;
    line-height: 1.75;
    color: var(--ink);
    -webkit-font-smoothing: antialiased;
}

/* ── Note Canvas ── */
.note-paper {
    background: var(--surface);
    border-radius: var(--radius);
    padding: 2.5rem 3rem;
    box-shadow: var(--shadow-lg);
    border: 1px solid var(--border);
    min-height: 500px;
    position: relative;
}
.note-paper::before {
    content: '';
    position: absolute;
    top: 0; left: 0; right: 0;
    height: 5px;
    background: linear-gradient(90deg, var(--indigo), var(--blue-mid), var(--teal));
    border-radius: var(--radius) var(--radius) 0 0;
}

/* ── Headings ── */
h1 {
    font-family: var(--fn);
    font-weight: 800;
    font-size: 1.9rem;
    color: var(--indigo);
    text-align: center;
    padding: 1.2rem 1.5rem;
    margin-bottom: 1.5rem;
    letter-spacing: -0.02em;
    background: linear-gradient(135deg, var(--indigo-lite) 0%, #dce8fd 100%);
    border-radius: var(--radius);
    border-left: 5px solid var(--indigo);
    box-shadow: var(--shadow-sm);
}
h2 {
    font-weight: 700;
    font-size: 1.25rem;
    color: var(--blue);
    margin-top: 2rem;
    margin-bottom: 0.6rem;
    padding-bottom: 0.4rem;
    border-bottom: 2px solid var(--blue-lite);
    display: flex;
    align-items: center;
    gap: 8px;
    letter-spacing: -0.01em;
}
h2::before {
    content: '';
    display: inline-block;
    width: 4px;
    height: 1.1em;
    background: linear-gradient(to bottom, var(--blue), var(--teal));
    border-radius: 2px;
    flex-shrink: 0;
}
h3 {
    font-weight: 600;
    font-size: 1.05rem;
    color: var(--emerald);
    margin-top: 1.2rem;
    margin-bottom: 0.4rem;
    padding-left: 0.6rem;
    border-left: 3px solid var(--emerald);
}

/* ── Chapter Map (pre block) ── */
pre {
    background: linear-gradient(135deg, var(--indigo-lite) 0%, #dce8fd 100%);
    border: 1px solid #c5cae9;
    border-left: 4px solid var(--indigo);
    border-radius: var(--radius-sm);
    padding: 1.25rem 1.5rem;
    font-family: var(--fn-mono);
    font-size: 0.88rem;
    overflow-x: auto;
    line-height: 1.8;
    color: var(--indigo);
    box-shadow: var(--shadow-sm);
}

/* ── Cards — Callout Boxes ── */
.highlight-box,
.activity-box,
.case-study-box,
.exam-box,
.memory-box,
.pause-ponder,
.threads-curiosity {
    border-radius: var(--radius-sm);
    padding: 1rem 1.2rem 1rem 1.5rem;
    margin: 1.1rem 0;
    box-shadow: var(--shadow-md);
    position: relative;
    border: 1px solid transparent;
}

/* Header label inside each card */
.highlight-box::before,
.activity-box::before,
.case-study-box::before,
.exam-box::before,
.memory-box::before,
.pause-ponder::before,
.threads-curiosity::before {
    font-family: var(--fn);
    font-weight: 700;
    font-size: 0.72rem;
    letter-spacing: 0.08em;
    text-transform: uppercase;
    display: block;
    margin-bottom: 0.5rem;
    padding-bottom: 0.35rem;
    border-bottom: 1px solid rgba(0,0,0,0.08);
}

/* Key Concept — Amber */
.highlight-box {
    background: linear-gradient(135deg, var(--amber-lite) 70%, #fff3cd);
    border-left: 5px solid var(--amber);
    border-color: #ffe082;
}
.highlight-box::before {
    content: "✏️  KEY CONCEPT";
    color: #e65100;
}

/* Activity / Experiment — Emerald */
.activity-box {
    background: linear-gradient(135deg, var(--emerald-lite) 70%, #dcedc8);
    border-left: 5px solid var(--emerald);
    border-color: #a5d6a7;
}
.activity-box::before {
    content: "🔬  EXPERIMENT / ACTIVITY";
    color: var(--emerald);
}

/* Real-World Case — Blue */
.case-study-box {
    background: linear-gradient(135deg, var(--blue-lite) 70%, #dce8fd);
    border-left: 5px solid var(--blue-mid);
    border-color: #90caf9;
}
.case-study-box::before {
    content: "🌍  REAL-WORLD APPLICATION";
    color: var(--blue);
}

/* Exam Corner — Amethyst */
.exam-box {
    background: linear-gradient(135deg, var(--amethyst-lite) 70%, #ede0f8);
    border-left: 5px solid var(--amethyst);
    border-color: #ce93d8;
}
.exam-box::before {
    content: "🎯  EXAM CORNER";
    color: var(--amethyst);
}

/* Memory Trick — Teal */
.memory-box {
    background: linear-gradient(135deg, var(--teal-lite) 70%, #b2dfdb);
    border-left: 5px solid var(--teal);
    border-color: #80cbc4;
    font-style: normal;
}
.memory-box::before {
    content: "🧠  MEMORY TRICK";
    color: var(--teal);
}

/* 🔍 Pause & Ponder — magnifying glass theme */
.pause-ponder {
    background: linear-gradient(135deg, #fce4ec 70%, #f8bbd9);
    border-left: 5px solid var(--rose);
    border-color: #ef9a9a;
}
.pause-ponder::before {
    content: "🔍  PAUSE & PONDER";
    color: var(--rose);
}

/* 🧭 Threads of Curiosity — compass theme */
.threads-curiosity {
    background: linear-gradient(135deg, var(--amethyst-lite) 60%, #e8eaf6);
    border-left: 5px solid var(--amethyst);
    border: 2px dashed #ce93d8;
    border-left: 5px solid var(--amethyst);
}
.threads-curiosity::before {
    content: "🧭  THREADS OF CURIOSITY";
    color: var(--amethyst);
}

/* ── Image Placeholder ── */
.img-placeholder {
    background: linear-gradient(135deg, var(--canvas2) 0%, #e8eaf6 100%);
    border: 2px dashed #9fa8da;
    border-radius: var(--radius);
    min-height: 180px;
    display: flex;
    flex-direction: column;
    align-items: center;
    justify-content: center;
    gap: 10px;
    margin: 1.1rem 0;
    padding: 1.5rem;
    text-align: center;
    color: var(--indigo-mid);
    font-size: 0.88rem;
    font-weight: 500;
    box-shadow: var(--shadow-sm);
    position: relative;
    overflow: hidden;
}
.img-placeholder::before {
    content: "🖼";
    font-size: 2.8rem;
    opacity: 0.6;
}
.img-placeholder::after {
    content: 'Scientific Illustration';
    position: absolute;
    bottom: 0; left: 0; right: 0;
    background: linear-gradient(90deg, var(--indigo), var(--blue-mid));
    color: white;
    font-size: 0.65rem;
    font-weight: 700;
    letter-spacing: 0.1em;
    text-transform: uppercase;
    padding: 4px 10px;
    text-align: center;
}

/* ── Sticky Note ── */
.sticky {
    background: linear-gradient(135deg, #fff9c4 80%, #fff176);
    border-radius: 4px 4px 20px 4px;
    padding: 0.85rem 1.1rem;
    margin: 1rem 0;
    box-shadow: 3px 4px 12px rgba(0,0,0,0.12), inset 0 -2px 4px rgba(0,0,0,0.04);
    transform: rotate(-0.6deg);
    display: inline-block;
    min-width: 200px;
    font-size: 0.93rem;
    font-weight: 500;
    color: #5d4037;
    border-top: 3px solid #fdd835;
}

/* ── Formula Block ── */
.formula-block {
    background: linear-gradient(135deg, #ede7f6 70%, #d1c4e9);
    border: 1.5px solid #b39ddb;
    border-left: 5px solid var(--amethyst);
    border-radius: var(--radius-sm);
    padding: 1.2rem 1.5rem;
    margin: 1.1rem 0;
    text-align: center;
    font-family: var(--fn-serif);
    font-size: 1.1rem;
    box-shadow: var(--shadow-md);
    overflow-x: auto;
}
.formula-block::before {
    content: "∑  FORMULA";
    display: block;
    font-family: var(--fn);
    font-size: 0.7rem;
    font-weight: 700;
    letter-spacing: 0.1em;
    text-transform: uppercase;
    color: var(--amethyst);
    margin-bottom: 0.5rem;
}

/* ── Grid Layouts ── */
.two-col {
    display: grid;
    grid-template-columns: 1fr 1fr;
    gap: 1rem;
    margin: 1rem 0;
}
.three-col {
    display: grid;
    grid-template-columns: 1fr 1fr 1fr;
    gap: 0.85rem;
    margin: 1rem 0;
}
.two-col > div,
.three-col > div {
    background: var(--canvas);
    border: 1px solid var(--border);
    border-radius: var(--radius-sm);
    padding: 0.85rem 1rem;
    font-size: 0.93rem;
}

/* ── Tables ── */
table {
    width: 100%;
    border-collapse: separate;
    border-spacing: 0;
    margin: 1.2rem 0;
    font-size: 0.92rem;
    border-radius: var(--radius-sm);
    overflow: hidden;
    box-shadow: var(--shadow-md);
    border: 1px solid var(--border);
}
th {
    background: linear-gradient(135deg, var(--indigo) 0%, var(--blue-mid) 100%);
    color: white;
    padding: 11px 16px;
    text-align: left;
    font-weight: 600;
    font-size: 0.85rem;
    letter-spacing: 0.03em;
    position: sticky;
    top: 0;
}
td {
    padding: 9px 16px;
    border-bottom: 1px solid var(--rule);
    background: var(--surface);
    vertical-align: top;
}
tr:nth-child(even) td { background: var(--canvas); }
tr:last-child td { border-bottom: none; }
tr:hover td { background: var(--blue-lite); transition: background 0.15s; }

/* ── Bullet List ── */
ul.bullet-list {
    padding-left: 0.5rem;
    margin: 0.6rem 0;
    list-style: none;
}
ul.bullet-list li {
    padding: 0.3rem 0 0.3rem 1.6rem;
    position: relative;
    border-bottom: 1px solid var(--rule);
    font-size: 0.95rem;
}
ul.bullet-list li:last-child { border-bottom: none; }
ul.bullet-list li::before {
    content: "◆";
    position: absolute;
    left: 0;
    color: var(--blue-mid);
    font-size: 0.65rem;
    top: 0.55rem;
}

/* ── Mobile Responsive ── */
@media (max-width: 768px) {
    .note-paper { padding: 1.5rem 1.2rem; }
    h1 { font-size: 1.4rem; padding: 1rem; }
    h2 { font-size: 1.1rem; }
    h3 { font-size: 1rem; }
    .two-col, .three-col { grid-template-columns: 1fr; gap: 0.75rem; }
    .sticky { transform: none; display: block; min-width: unset; width: 100%; }
    table { font-size: 0.82rem; }
    th, td { padding: 7px 10px; }
    pre { font-size: 0.8rem; padding: 1rem; }
    .formula-block { font-size: 0.95rem; padding: 1rem; }
    .highlight-box, .activity-box, .case-study-box,
    .exam-box, .memory-box, .pause-ponder, .threads-curiosity { padding: 0.85rem 1rem; }
    .img-placeholder { min-height: 130px; }
}
@media (max-width: 480px) {
    body { font-size: 14px; }
    .note-paper { padding: 1rem 0.85rem; }
    h1 { font-size: 1.2rem; }
}
"""

# ── HTML wrapper for exported ZIP files ──────────────────────────────────────
HTML_WRAPPER = """<!DOCTYPE html>
<html lang="en">
<head>
  <meta charset="UTF-8">
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  <title>{title}</title>
  <link rel="preconnect" href="https://fonts.googleapis.com">
  <link href="https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800&family=Lora:ital,wght@0,400;0,600;1,400&family=JetBrains+Mono:wght@400;500&display=swap" rel="stylesheet">
  <link rel="stylesheet" href="style.css">
  <link rel="stylesheet" href="https://cdn.jsdelivr.net/npm/katex@0.16.9/dist/katex.min.css">
  <script defer src="https://cdn.jsdelivr.net/npm/katex@0.16.9/dist/katex.min.js"></script>
  <script defer src="https://cdn.jsdelivr.net/npm/katex@0.16.9/dist/contrib/auto-render.min.js"
          onload="renderMathInElement(document.body,{{delimiters:[{{left:'$$',right:'$$',display:true}},{{left:'$',right:'$',display:false}}]}});"></script>
</head>
<body>
  <div class="note-paper">
    {body}
  </div>
</body>
</html>"""

# ═══════════════════════════════════════════════════════════════════════════════
#  EXTRACTION HELPERS
# ═══════════════════════════════════════════════════════════════════════════════
def extract_pdf(buf: io.BytesIO) -> str:
    reader = pypdf.PdfReader(buf)
    return "\n\n".join(p.extract_text() for p in reader.pages if p.extract_text())


def extract_pptx(buf: io.BytesIO) -> str:
    prs = Presentation(buf)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        parts = [f"<h3>Slide {i}</h3>"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        slides.append("\n".join(parts))
    return "\n\n".join(slides)


def sanitise(html: str) -> str:
    html = re.sub(r"\*\*(.+?)\*\*", r"<b>\1</b>", html)
    html = re.sub(r"\*(.+?)\*", r"<em>\1</em>", html)
    html = re.sub(r"(?m)^[ \t]*[-•]\s+(.+)$", r"<li>\1</li>", html)
    html = re.sub(r"((?:<li>.*?</li>\s*)+)", r'<ul class="bullet-list">\1</ul>', html, flags=re.S)
    html = re.sub(r"```[a-z]*\n?", "", html)
    return html


# ═══════════════════════════════════════════════════════════════════════════════
#  PROMPT BUILDER
# ═══════════════════════════════════════════════════════════════════════════════
def build_prompt(raw_text: str, subject: str) -> str:
    return f"""You are an expert Science educator, NCERT content analyst, and topper-level student note-maker for subject "{subject}".

Convert the provided chapter into premium structured notes for a modern digital notebook with a scientific aesthetic.

OUTPUT RULES — follow exactly:
1. Return ONLY raw HTML. No markdown code fences. No preamble or commentary.
2. Use exactly these three section wrappers:
   <div id="visual-notes">   ... </div>
   <div id="summary-page">   ... </div>
   <div id="formula-matrix"> ... </div>

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
SECTION 1 — <div id="visual-notes">
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
Generate COMPLETE TOPPER NOTES covering EVERY topic.

⚠️ MANDATORY COVERAGE:
• Go through source material LINE BY LINE — nothing skipped.
• Every heading, subheading, example, activity, table, figure caption MUST appear.
• Write MORE not less. Completeness is #1 priority.

A) CHAPTER MAP — use <pre> tags with tree structure:
   CHAPTER TITLE
   │
   ├── 1. Topic One
   │    ├── 1.1 Subtopic
   │    └── 1.2 Subtopic
   └── N. Topic N

B) TOPIC-WISE NOTES for every topic:
   • <h2> for main topics · <h3> for subtopics
   • Key definitions/laws       → <div class="highlight-box">
   • Activities/experiments     → <div class="activity-box">
   • Real-world applications    → <div class="case-study-box">
   • Exam tips                  → <div class="exam-box">
   • Memory tricks/mnemonics    → <div class="memory-box">
   • Evidence/observation Qs    → <div class="pause-ponder">
   • "Threads of Curiosity"     → <div class="threads-curiosity">
   • Bullet points              → <ul class="bullet-list"><li>
   • Inline emoji: ⭐ Important  🔥 Very Important  ⚠️ Common Mistake  🎯 Exam  📝 Remember  💡 Tip

   LAYOUTS:
   • Comparisons: <div class="two-col"><div>...</div><div>...</div></div>
   • Quick facts: <div class="three-col"><div>...</div>...</div>
   • Tips:        <div class="sticky">short tip</div>
   • Diagrams:    <div class="img-placeholder">Diagram: [describe it precisely]</div>

C) SCIENTIFIC TERMS — in a highlight-box, one per line: Term → Meaning
D) EXAMPLES — 🎯 Example N — Situation | Variables | Learning
E) ACTIVITIES — 🧪 Activity — Aim | Observation | Conclusion | Exam Q
F) MEMORY TOOLS — mnemonics, recall tricks, last-minute tips

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
SECTION 2 — <div id="summary-page">
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
ONE-PAGE REVISION SHEET — ALL definitions, laws, examples, common mistakes, 10 exam points.
Use <h2> for categories and <ul class="bullet-list"><li> for all items.

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
SECTION 3 — <div id="formula-matrix">
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
Every formula/equation/constant → <div class="formula-block">
Include: LaTeX formula ($$...$$) · symbol meanings · units · conditions · exam tip.
If no formulas, make "📌 Important Relationships" section using formula-block style.

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
SOURCE MATERIAL (subject: {subject}):
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
{raw_text}
"""


# ═══════════════════════════════════════════════════════════════════════════════
#  SECTION PARSER
# ═══════════════════════════════════════════════════════════════════════════════
def slice_div(div_id: str, source: str) -> str:
    open_tag = f'<div id="{div_id}">'
    start = source.find(open_tag)
    if start == -1:
        return f"<p><em>Section '{div_id}' not found. Try re-processing.</em></p>"
    start += len(open_tag)
    depth, i = 1, start
    while i < len(source) and depth > 0:
        op = source.find("<div", i)
        cl = source.find("</div>", i)
        if cl == -1:
            break
        if op != -1 and op < cl:
            depth += 1; i = op + 4
        else:
            depth -= 1
            if depth == 0:
                return sanitise(source[start:cl].strip())
            i = cl + 6
    return sanitise(source[start:].strip())


def parse_sections(raw: str) -> dict:
    return {
        "notes":    slice_div("visual-notes",   raw),
        "bullets":  slice_div("summary-page",   raw),
        "formulas": slice_div("formula-matrix", raw),
    }


# ═══════════════════════════════════════════════════════════════════════════════
#  API CALLS
# ═══════════════════════════════════════════════════════════════════════════════
# Model name map for display
MODEL_LABELS = {
    "claude-haiku-4-5-20251001": "Claude Haiku 4.5",
    "claude-sonnet-4-6":         "Claude Sonnet 4.6",
    "claude-opus-4-8":           "Claude Opus 4.8",
    "gemini-2.0-flash-lite":     "Gemini 2.0 Flash Lite",
    "gemini-1.5-flash":          "Gemini 1.5 Flash",
    "gemini-2.0-flash":          "Gemini 2.0 Flash",
    "gpt-4o-mini":               "GPT-4o Mini",
    "gpt-4o":                    "GPT-4o",
    "gpt-4.5-preview":           "GPT-4.5",
    "gpt-5.2":                   "GPT-5.2",
}


def call_claude(raw_text: str, api_keys: list, subject: str, model: str) -> dict:
    prompt = build_prompt(raw_text, subject)
    try:
        import anthropic
    except ImportError:
        st.error("Run: `pip install anthropic`")
        st.stop()
    last_err = None
    for key in api_keys:
        try:
            client = anthropic.Anthropic(api_key=key)
            resp = client.messages.create(
                model=model,
                max_tokens=16000,
                messages=[{"role": "user", "content": prompt}],
            )
            return parse_sections(resp.content[0].text.strip())
        except Exception as e:
            last_err = e
    raise last_err


def call_openai(raw_text: str, api_keys: list, subject: str, model: str) -> dict:
    prompt = build_prompt(raw_text, subject)
    try:
        from openai import OpenAI
    except ImportError:
        st.error("Run: `pip install openai`")
        st.stop()
    last_err = None
    for key in api_keys:
        try:
            client = OpenAI(api_key=key)
            resp = client.chat.completions.create(
                model=model,
                max_tokens=16000,
                messages=[{"role": "user", "content": prompt}],
            )
            return parse_sections(resp.choices[0].message.content.strip())
        except Exception as e:
            last_err = e
    raise last_err


def call_gemini(raw_text: str, api_keys: list, subject: str, model: str) -> dict:
    prompt = build_prompt(raw_text, subject)
    last_err = None
    for key in api_keys:
        try:
            try:
                from google import genai as _g
                from google.genai import types as gt
                client = _g.Client(api_key=key)
                cfg = gt.GenerateContentConfig(max_output_tokens=16000)
                resp = client.models.generate_content(model=model, contents=prompt, config=cfg)
            except ImportError:
                import google.generativeai as genai_legacy
                genai_legacy.configure(api_key=key)
                resp = genai_legacy.GenerativeModel(model).generate_content(
                    prompt, generation_config={"max_output_tokens": 16000})
            return parse_sections(resp.text.strip())
        except Exception as e:
            last_err = e
    raise last_err


def call_ollama(raw_text: str, model: str, subject: str) -> dict:
    import urllib.request, json
    prompt = build_prompt(raw_text, subject)
    payload = json.dumps({
        "model": model, "prompt": prompt, "stream": False,
        "options": {"num_predict": 6000, "num_ctx": 16384}
    }).encode()
    req = urllib.request.Request(
        "http://localhost:11434/api/generate",
        data=payload,
        headers={"Content-Type": "application/json"},
    )
    try:
        with urllib.request.urlopen(req, timeout=900) as resp:
            raw = json.loads(resp.read())["response"].strip()
    except Exception as e:
        raise RuntimeError(
            f"Ollama error: {e}. Ensure Ollama is running (`ollama serve`) "
            f"and model '{model}' is pulled (`ollama pull {model}`)."
        )
    return parse_sections(raw)


# ── ZIP builder ────────────────────────────────────────────────────────────────
def build_zip(data: dict, base_name: str) -> io.BytesIO:
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("style.css",      CSS)
        zf.writestr("index.html",     HTML_WRAPPER.format(title="Detailed Notes",   body=data["notes"]))
        zf.writestr("summary.html",   HTML_WRAPPER.format(title="Quick Highlights", body=data["bullets"]))
        zf.writestr("formulas.html",  HTML_WRAPPER.format(title="Formula Sheet",    body=data["formulas"]))
        zf.writestr("README.txt",
            "Open index.html, summary.html, or formulas.html in any browser.\n"
            "KaTeX formulas render when connected to the internet.\n"
            f"Generated from: {base_name}\n")
    buf.seek(0)
    return buf


# ═══════════════════════════════════════════════════════════════════════════════
#  STREAMLIT UI
# ═══════════════════════════════════════════════════════════════════════════════
st.set_page_config(page_title="Production Note Studio", page_icon="🔬", layout="wide")
st.markdown(f"<style>{CSS}</style>", unsafe_allow_html=True)

st.title("🔬 Production Note Studio")
st.caption("Transform textbook PDFs & lecture decks into premium scientific notes.")

# ── Sidebar ────────────────────────────────────────────────────────────────────
with st.sidebar:
    st.markdown("## ⚙️ Configuration")
    st.divider()

    # ── Provider ──
    st.markdown("### 🤖 AI Provider")
    provider = st.radio(
        "Provider",
        ["🏠 Ollama", "🆓 Gemini", "💎 Claude", "🧠 OpenAI"],
        label_visibility="collapsed",
        help="Ollama = free local · Gemini = free tier · Claude/OpenAI = paid"
    )

    use_ollama = provider.startswith("🏠")
    use_gemini = provider.startswith("🆓")
    use_openai = provider.startswith("🧠")
    use_claude = provider.startswith("💎")

    # ── Model variant ──
    st.markdown("### 🔧 Model")

    if use_ollama:
        ollama_model = st.selectbox(
            "Local model",
            ["gemma2:2b", "llama3.2", "mistral", "phi3", "qwen2.5", "llava"],
            help="Run `ollama pull <model>` once to download."
        )
        selected_model = ollama_model
        api_keys = []
        st.info("Make sure Ollama is running:\n```\nollama serve\n```")

    elif use_gemini:
        selected_model = st.selectbox(
            "Gemini model",
            ["gemini-2.0-flash-lite", "gemini-1.5-flash", "gemini-2.0-flash"],
            help="Flash Lite = fastest & free · 2.0 Flash = best quality"
        )
        key_input = st.text_input(
            "Gemini API Key", type="password", placeholder="AIza…",
            help="[Get free key →](https://aistudio.google.com/apikey)"
        )
        api_keys = [key_input] if key_input.strip() else GEMINI_KEYS
        st.caption("[Get free Gemini key →](https://aistudio.google.com/apikey)")

    elif use_claude:
        selected_model = st.selectbox(
            "Claude model",
            ["claude-haiku-4-5-20251001", "claude-sonnet-4-6", "claude-opus-4-8"],
            format_func=lambda m: MODEL_LABELS.get(m, m),
            help="Haiku = cheapest · Sonnet = balanced · Opus = best quality"
        )
        key_input = st.text_input(
            "Claude API Key", type="password", placeholder="sk-ant-…",
            help="[Get key →](https://console.anthropic.com/)"
        )
        api_keys = [key_input] if key_input.strip() else CLAUDE_KEYS
        st.caption("[Get Claude key →](https://console.anthropic.com/)")

    else:  # OpenAI
        selected_model = st.selectbox(
            "OpenAI model",
            ["gpt-4o-mini", "gpt-4o", "gpt-4.5-preview", "gpt-5.2"],
            format_func=lambda m: MODEL_LABELS.get(m, m),
            help="4o-mini = cheapest · 4o = best · 4.5 = most powerful"
        )
        key_input = st.text_input(
            "OpenAI API Key", type="password", placeholder="Paste your OpenAI key…",
            help="[Get key →](https://platform.openai.com/api-keys)"
        )
        api_keys = [key_input] if key_input.strip() else OPENAI_KEYS
        st.caption("[Get OpenAI key →](https://platform.openai.com/api-keys)")

    # ── Subject ──
    st.divider()
    st.markdown("### 📚 Subject")
    subject = st.selectbox(
        "Subject",
        ["Science 🔬", "Biology 🧬", "Physics ⚡", "Chemistry 🧪",
         "Mathematics 📐", "History 📜", "Geography 🌍",
         "Economics 📊", "Computer Science 💻"],
        label_visibility="collapsed",
    )

    # ── Info footer ──
    st.divider()
    model_display = MODEL_LABELS.get(selected_model, selected_model) if not use_ollama else selected_model
    st.markdown(f"**Active model:** `{model_display}`")
    st.markdown("**Formats:** PDF · PPTX")
    if use_ollama:
        st.markdown("[Install Ollama →](https://ollama.com/download)")


# ── Main area ─────────────────────────────────────────────────────────────────
uploaded = st.file_uploader(
    "Upload a Chapter PDF or Lecture Deck (.pptx)",
    type=["pdf", "pptx"]
)

if uploaded:
    if not use_ollama and not api_keys:
        provider_name = "Gemini" if use_gemini else ("OpenAI" if use_openai else "Anthropic")
        var_name = "GEMINI_API_KEYS" if use_gemini else ("OPENAI_API_KEYS" if use_openai else "CLAUDE_API_KEYS")
        st.error(
            f"No {provider_name} API key found. "
            f"Enter it in the sidebar or set the `{var_name}` environment variable."
        )
        st.stop()

    ext = uploaded.name.rsplit(".", 1)[-1].lower()
    base_name = uploaded.name
    model_display = MODEL_LABELS.get(selected_model, selected_model) if not use_ollama else selected_model

    col1, col2 = st.columns([3, 1])
    with col1:
        st.info(f"📄 **{base_name}** — ready to process with **{model_display}**")
    with col2:
        run_btn = st.button("⚡ Generate Notes", type="primary", use_container_width=True)

    if run_btn:
        with st.status("Processing your document…", expanded=True) as status:
            st.write("📖 Extracting text from file…")
            file_buf = io.BytesIO(uploaded.read())

            raw_text = extract_pdf(file_buf) if ext == "pdf" else extract_pptx(file_buf)
            word_count = len(raw_text.split())
            st.write(f"✅ Extracted **{word_count:,} words** from {base_name}")

            if word_count < 50:
                status.update(label="Extraction failed — too little text found.", state="error")
                st.error("The file has very little extractable text. Scanned image PDFs are not supported.")
                st.stop()

            MAX_WORDS = 4000 if use_ollama else (5000 if use_gemini else 12000)
            if word_count > MAX_WORDS:
                raw_text = " ".join(raw_text.split()[:MAX_WORDS])
                st.warning(f"⚠️ Truncated to first {MAX_WORDS:,} words (chapter had {word_count:,}). Use Claude/GPT for full coverage.")

            st.write(f"🤖 Sending to **{model_display}** for note compilation…")
            try:
                if use_ollama:
                    sections = call_ollama(raw_text, selected_model, subject)
                elif use_gemini:
                    sections = call_gemini(raw_text, api_keys, subject, selected_model)
                elif use_openai:
                    sections = call_openai(raw_text, api_keys, subject, selected_model)
                else:
                    sections = call_claude(raw_text, api_keys, subject, selected_model)
            except Exception as e:
                status.update(label="Error — see details below.", state="error")
                st.error(f"**{model_display} error:** {e}")
                if use_ollama:
                    st.info("Ensure Ollama is running: `ollama serve` · Pull model: `ollama pull " + selected_model + "`")
                elif use_gemini:
                    st.info("pip install google-genai · Get free key: https://aistudio.google.com/apikey")
                elif use_openai:
                    st.info("pip install openai · Get key: https://platform.openai.com/api-keys")
                else:
                    st.info("pip install anthropic · Get key: https://console.anthropic.com/")
                st.stop()

            st.write("📦 Building ZIP export package…")
            zip_buf = build_zip(sections, base_name)
            status.update(label="✅ Notes ready!", state="complete", expanded=False)

        st.session_state["sections"]  = sections
        st.session_state["zip_buf"]   = zip_buf
        st.session_state["base_name"] = base_name

# ── Results ───────────────────────────────────────────────────────────────────
if "sections" in st.session_state:
    sections  = st.session_state["sections"]
    zip_buf   = st.session_state["zip_buf"]
    base_name = st.session_state["base_name"]

    st.success(f"✅ Notebook generated for **{base_name}**")

    tab1, tab2, tab3 = st.tabs(["🗒️ Detailed Notes", "⚡ Quick Highlights", "∑ Formula Sheet"])

    with tab1:
        st.markdown('<div class="note-paper">', unsafe_allow_html=True)
        st.markdown(sections["notes"], unsafe_allow_html=True)
        st.markdown('</div>', unsafe_allow_html=True)

    with tab2:
        st.markdown('<div class="note-paper">', unsafe_allow_html=True)
        st.markdown(sections["bullets"], unsafe_allow_html=True)
        st.markdown('</div>', unsafe_allow_html=True)

    with tab3:
        st.markdown('<div class="note-paper">', unsafe_allow_html=True)
        st.markdown(sections["formulas"], unsafe_allow_html=True)
        st.markdown('</div>', unsafe_allow_html=True)

    st.divider()
    st.download_button(
        label="⬇️ Download Offline HTML + CSS Bundle (.zip)",
        data=zip_buf,
        file_name=f"{base_name.rsplit('.', 1)[0]}_notes.zip",
        mime="application/zip",
        type="primary",
    )
